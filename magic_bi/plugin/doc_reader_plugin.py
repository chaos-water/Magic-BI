# from pydantic import BaseModel
import io
import pptx
import docx
from loguru import logger
# import fitz


# class PdfParser():
#     def parse(self, file_name: str, file_bytes: bytes) -> str:
#         text_data = ""
#         doc = fitz.open(stream=file_bytes)
#         page_count = doc.page_count
#
#         for page_num in range(page_count):
#             text = doc.get_page_text(page_num)
#
#             text = text.strip(" ")
#             if len(text) > 0:
#                 text_data += text + "\f"
#
#         text_data = self.file_clean_text(text_data)
#         logger.debug("parse %s suc, txt_len:%d" % (file_name, len(text_data)))
#         return text_data


class PptxParser():
    def parse(self, file_name: str, file_bytes: bytes) -> str:
        text_data = ""
        try:
            data = io.BytesIO(file_bytes)
            ppt = pptx.Presentation(data)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_data += shape.text

            logger.debug("parse %s suc, txt_len:%d" % (file_name, len(text_data)))
            return text_data
        except Exception as error:
            logger.error("catch exception: %s", str(error))
            return ""
        finally:
            data.close()


class DocxParser():
    def parse(self, file_name: str, file_bytes: bytes) -> str:
        text_data = ""
        try:
            data = io.BytesIO(file_bytes)
            doc = docx.Document(data)
            for para in doc.paragraphs:
                text_data += para.text

            logger.debug("parse %s suc, txt_len:%d" % (file_name, len(text_data)))
            return text_data
        except Exception as error:
            logger.error("catch exception: %s" % str(error))
            return ""
        finally:
            data.close()


class DocReaderPlugin():
    def __init__(self):
    # _pdf_parser: PdfParser = PdfParser()
        self._docx_parser: DocxParser = DocxParser()
        self._pptx_parser: PptxParser = PptxParser()

    def run(self, file_path: str) -> str:
        # file_path = file_path_or_url
        # if file_path_or_url.startswith("http"):
        #     pass
        #     file_path = self._download_file_bytes(file_path_or_url)

        if file_path.endswith(".pdf"):
            text_data = self._pdf_parser.parse(file_path)
        elif file_path.endswith(".docx"):
            text_data = self._docx_parser.parse(file_path)
        elif file_path.endswith(".pptx"):
            text_data = self._pptx_parser.parse(file_path)
        else:
            return ""

        return text_data

    def _download_file_bytes(self, file_url: str) -> str:
        pass

if __name__ == "__main__":
    docx_parser = DocxParser()
    file_path = "/Users/luguanglong/Codes/mojing/Magic-Data/DSL_modem设置操作和故障排除手册.docx"
    with open(file_path, "rb") as f:
        file_bytes = f.read()
        docx_parser.parse("test", file_bytes)
